import textwrap

import openai
from pptx import Presentation
from transformers import pipeline
from pprint import pprint
from PIL import Image
import os
import pytesseract

pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'


def extract_text_from_pptx(file_path, image_output_dir="images"):
    """
    Extracts text from a PowerPoint (.pptx) file.

    Args:
        file_path (str): The path to the PowerPoint file.

    Returns:
        List[str]: A list of text content for each slide.
    """
    pres = Presentation(file_path)
    text_content = []
    if not os.path.exists(image_output_dir):
        os.makedirs(image_output_dir)

    for i, slide in enumerate(pres.slides):
        slide_text = []

        # Extract text from shapes
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text.append(shape.text)

        # Extract text from images
        for shape in slide.shapes:
            if shape.shape_type == 13:  # 13 corresponds to a picture
                image = shape.image
                image_path = os.path.join(image_output_dir, f"slide_{i + 1}_image.png")
                with open(image_path, "wb") as img_file:
                    img_file.write(image.blob)

                # Perform OCR on the saved image
                ocr_text = pytesseract.image_to_string(Image.open(image_path))
                slide_text.append(f"Image OCR: {ocr_text}")

        # Combine text and OCR results for the slide
        text_content.append(" ".join(slide_text))

    return text_content


# Example usage:
file_path = r"COP 4331 - Gathering Requirements.pptx" 
text_data = extract_text_from_pptx(file_path)
print(text_data)

# Load a pre-trained summarization model
summarizer = pipeline("summarization", model="t5-small")


def summarize_slides(text_data, max_length=40):
    """
    Summarizes the text content of each slide.

    Args:
        text_data (List[str]): List of text content for each slide.
        max_length (int): The maximum length of the summary.

    Returns:
        List[str]: Summarized content for each slide.
    """
    summaries = []
    for text in text_data:
        summary = summarizer(text, max_length=max_length, min_length=30, do_sample=False)[0]['summary_text']
        summaries.append(summary)
    return summaries


def chunk_text(text, chunk_size=512):
    """Split text into chunks that fit within the model's token limit."""
    return textwrap.wrap(text, chunk_size)


# Load a pre-trained text generation model for question generation
question_generator = pipeline("text2text-generation", model="t5-small")


def generate_quiz_questions(text_data, max_length=50):
    """
    Generates quiz questions from the entire text content of all slides using OpenAI GPT.

    Args:
        text_data (List[str]): List of summarized text content for each slide.
        max_length (int): The maximum length for generated questions.

    Returns:
        List[str]: Quiz questions generated from the entire content.
    """
    # Combine all the text data into a single string
    combined_text = " ".join(text_data)

    # Create a prompt for GPT to generate a quiz based on the combined text
    prompt = f"Based on the following summarized content, generate a list of quiz questions:\n\n{combined_text}\n\nPlease make the questions varied and engaging, and ensure they are suitable for a quiz."

    '''
    # Call the OpenAI API to generate the quiz questions
    response = openai.completions.create(
        model="gpt-3.5",  # Use 'gpt-4' if you have access
        prompt=prompt,
        max_tokens=1000,  # Adjust the max tokens as needed
        temperature=0.7,  # Adjust the temperature for creativity (higher means more creative)
        n=1,  # Generate one completion
    )'''
    response = openai.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt}
        ]
    )

    # Get the generated questions from the response
    generated_text = response.choices[0].message.content

    # Split the generated text into individual questions, assuming each question is separated by a newline
    questions = generated_text.split("\n")

    return questions


# Example usage:
summaries = summarize_slides(text_data)
print(summaries)

# Generate quiz questions from summaries
quiz_questions = generate_quiz_questions(summaries)
for question in quiz_questions:
    print(question)
